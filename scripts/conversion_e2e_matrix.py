from __future__ import annotations

import argparse
import html
import json
import re
import shutil
import subprocess
import sys
import tempfile
import zipfile
from pathlib import Path

import fitz
from morphix_worker.converters.registry import LocalConverterRegistry


MARKER = "MORPHIX_TEST_MARKER_20260719"
EXPECTED_DURATION_SECONDS = 2.0

# Intentionally independent from the runtime registry. A change to production
# cannot make this matrix pass without the corresponding end-to-end behavior.
MATRIX: tuple[tuple[str, str], ...] = (
    ("docx", "pdf"), ("xlsx", "pdf"), ("pptx", "pdf"), ("pdf", "docx"),
    ("pdf", "txt"), ("pdf", "html"), ("pdf", "png"), ("pdf", "jpg"),
    ("docx", "txt"), ("docx", "odt"), ("docx", "html"), ("xlsx", "ods"),
    ("pptx", "odp"), ("pptx", "txt"), ("csv", "xlsx"), ("xlsx", "csv"),
    ("csv", "pdf"), ("csv", "ods"), ("png", "jpg"), ("jpg", "png"),
    ("jpg", "webp"), ("png", "webp"), ("png", "gif"), ("png", "tiff"),
    ("png", "bmp"), ("jpg", "gif"), ("jpg", "tiff"), ("jpg", "bmp"),
    ("svg", "png"), ("svg", "jpg"), ("gif", "png"), ("tiff", "jpg"),
    ("mp4", "mp3"), ("mp4", "webm"), ("mov", "mp4"), ("avi", "mp4"),
    ("avi", "webm"), ("mkv", "mp4"), ("mkv", "webm"), ("mov", "webm"),
    ("m4v", "mp4"), ("flv", "mp4"), ("3gp", "mp4"), ("wav", "mp3"),
    ("mp3", "wav"), ("aac", "mp3"), ("flac", "mp3"), ("ogg", "mp3"),
    ("m4a", "mp3"), ("opus", "wav"), ("aiff", "mp3"), ("aiff", "wav"),
)


def command(args: list[str], *, capture: bool = True) -> subprocess.CompletedProcess[str]:
    try:
        return subprocess.run(args, check=True, text=True, capture_output=capture)
    except subprocess.CalledProcessError as error:
        details = (error.stderr or error.stdout or "").strip()
        raise RuntimeError(f"Command failed ({error.returncode}): {' '.join(args)}\n{details}") from error


def binary(name: str) -> str:
    path = shutil.which(name)
    if not path:
        raise RuntimeError(f"Required binary not found: {name}")
    return path


def generate_fixtures(root: Path) -> dict[str, Path]:
    fixture_dir = root / "fixtures"
    fixture_dir.mkdir()
    fixtures: dict[str, Path] = {}
    ffmpeg = binary("ffmpeg")
    imagemagick = shutil.which("magick") or binary("convert")

    from docx import Document

    document = Document()
    document.add_heading(MARKER, level=1)
    document.add_paragraph("Conversion quality fixture")
    docx = fixture_dir / "marker.docx"
    document.save(docx)
    fixtures["docx"] = docx

    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet["A1"] = MARKER
    sheet["A2"] = "Conversion quality fixture"
    sheet["B2"] = 20260719
    xlsx = fixture_dir / "marker.xlsx"
    workbook.save(xlsx)
    fixtures["xlsx"] = xlsx

    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    textbox.text_frame.text = MARKER
    pptx = fixture_dir / "marker.pptx"
    presentation.save(pptx)
    fixtures["pptx"] = pptx

    csv = fixture_dir / "marker.csv"
    csv.write_text(f"name,value\n{MARKER},20260719\n", encoding="utf-8")
    fixtures["csv"] = csv

    pdf = fixture_dir / "marker.pdf"
    document = fitz.open()
    for page_number in (1, 2):
        page = document.new_page(width=595, height=842)
        page.insert_text((72, 90), f"{MARKER} PAGE {page_number}")
        page.insert_text((72, 120), "Conversion quality fixture")
    document.save(str(pdf))
    document.close()
    fixtures["pdf"] = pdf

    svg = fixture_dir / "marker.svg"
    svg.write_text(
        f'<svg xmlns="http://www.w3.org/2000/svg" width="320" height="180"><rect width="320" height="180" fill="#3047D8"/><circle cx="160" cy="90" r="55" fill="#F25A38"/><text x="160" y="95" text-anchor="middle" fill="white">{MARKER}</text></svg>',
        encoding="utf-8",
    )
    fixtures["svg"] = svg
    png = fixture_dir / "marker.png"
    command([imagemagick, "-size", "320x180", "gradient:#3047D8-#F25A38", str(png)])
    fixtures["png"] = png
    for extension in ("jpg", "gif", "tiff"):
        output = fixture_dir / f"marker.{extension}"
        command([imagemagick, str(png), str(output)])
        fixtures[extension] = output

    video_format_args = {
        "mp4": ["-c:v", "libx264", "-c:a", "aac", "-pix_fmt", "yuv420p"],
        "mov": ["-c:v", "libx264", "-c:a", "aac", "-pix_fmt", "yuv420p"],
        "avi": ["-c:v", "mpeg4", "-c:a", "mp3"],
        "mkv": ["-c:v", "libx264", "-c:a", "aac", "-pix_fmt", "yuv420p"],
        "m4v": ["-c:v", "libx264", "-c:a", "aac", "-pix_fmt", "yuv420p"],
        "flv": ["-c:v", "flv", "-c:a", "mp3"],
        "3gp": ["-c:v", "h263", "-c:a", "aac", "-ar", "8000"],
    }
    for extension, codec_args in video_format_args.items():
        output = fixture_dir / f"marker.{extension}"
        video_size = "176x144" if extension == "3gp" else "320x180"
        command([
            ffmpeg, "-y", "-f", "lavfi", "-i", f"color=c=#3047D8:s={video_size}:r=24",
            "-f", "lavfi", "-i", "sine=frequency=440:sample_rate=44100", "-t", str(EXPECTED_DURATION_SECONDS),
            *codec_args, "-shortest", str(output),
        ])
        fixtures[extension] = output

    audio_format_args = {
        "wav": ["-c:a", "pcm_s16le"], "mp3": ["-c:a", "libmp3lame"],
        "aac": ["-c:a", "aac"], "flac": ["-c:a", "flac"],
        "ogg": ["-c:a", "libvorbis"], "m4a": ["-c:a", "aac"],
        "opus": ["-c:a", "libopus"], "aiff": ["-c:a", "pcm_s16be"],
    }
    for extension, codec_args in audio_format_args.items():
        output = fixture_dir / f"marker.{extension}"
        command([
            ffmpeg, "-y", "-f", "lavfi", "-i", "sine=frequency=440:sample_rate=44100",
            "-t", str(EXPECTED_DURATION_SECONDS), *codec_args, str(output),
        ])
        fixtures[extension] = output

    missing = sorted({source for source, _ in MATRIX} - fixtures.keys())
    if missing:
        raise RuntimeError(f"Missing fixtures: {', '.join(missing)}")
    return fixtures


def extract_text(path: Path) -> str:
    if path.suffix.lower() == ".pdf":
        document = fitz.open(str(path))
        try:
            return "\n".join(page.get_text() for page in document)
        finally:
            document.close()
    if path.suffix.lower() in {".txt", ".csv", ".html"}:
        raw = path.read_text(encoding="utf-8", errors="ignore")
        return re.sub(r"<[^>]+>", " ", html.unescape(raw))
    if zipfile.is_zipfile(path):
        chunks: list[str] = []
        with zipfile.ZipFile(path) as archive:
            for name in archive.namelist():
                if name.endswith((".xml", ".rels")):
                    chunks.append(re.sub(r"<[^>]+>", " ", archive.read(name).decode("utf-8", errors="ignore")))
        return " ".join(chunks)
    return ""


def image_dimensions(path: Path, imagemagick: str) -> tuple[int, int]:
    result = command([imagemagick, "identify", "-format", "%w %h", str(path)]) if Path(imagemagick).name == "magick" else command([binary("identify"), "-format", "%w %h", str(path)])
    width, height = result.stdout.strip().split()[:2]
    return int(width), int(height)


def media_info(path: Path, ffprobe: str) -> dict[str, object]:
    result = command([ffprobe, "-v", "error", "-show_entries", "format=duration:stream=codec_type", "-of", "json", str(path)])
    return json.loads(result.stdout)


def normalize_image(path: Path, normalized_path: Path, imagemagick: str) -> None:
    if path.suffix.lower() == ".svg":
        rsvg = binary("rsvg-convert")
        command([rsvg, "--output", str(normalized_path), str(path)])
        return
    command([imagemagick, str(path), str(normalized_path)])


def visual_quality(source: Path, output: Path, source_format: str, target_format: str, pair_dir: Path, imagemagick: str) -> tuple[float, float]:
    normalized_source = pair_dir / ".source-quality.png"
    normalized_output = pair_dir / ".output-quality.png"
    if source_format == "pdf":
        document = fitz.open(str(source))
        try:
            document[0].get_pixmap(dpi=150, alpha=False).save(str(normalized_source))
        finally:
            document.close()
    else:
        normalize_image(source, normalized_source, imagemagick)
    normalize_image(output, normalized_output, imagemagick)

    compare = binary("compare")
    ssim_result = subprocess.run(
        [compare, "-metric", "SSIM", str(normalized_source), str(normalized_output), "null:"],
        text=True,
        capture_output=True,
    )
    psnr_result = subprocess.run(
        [compare, "-metric", "PSNR", str(normalized_source), str(normalized_output), "null:"],
        text=True,
        capture_output=True,
    )
    ssim_text = f"{ssim_result.stdout} {ssim_result.stderr}"
    psnr_text = f"{psnr_result.stdout} {psnr_result.stderr}"
    ssim_match = re.search(r"(?<![\d.])(1(?:\.0+)?|0?\.\d+)(?=\s*(?:\(|$))", ssim_text)
    psnr_match = re.search(r"(?<![\d.])(?:inf|\d+(?:\.\d+)?)(?=\s*(?:\(|$))", psnr_text, flags=re.IGNORECASE)
    if not ssim_match or not psnr_match:
        raise RuntimeError(f"Could not parse image quality metrics: SSIM={ssim_text!r} PSNR={psnr_text!r}")
    ssim = float(ssim_match.group(1))
    psnr = float("inf") if psnr_match.group(0).lower() == "inf" else float(psnr_match.group(0))
    threshold = 0.85 if source_format == "svg" else (0.90 if target_format in {"jpg", "webp", "gif"} else 0.99)
    if ssim < threshold:
        raise ValueError(f"SSIM {ssim:.4f} is below the {threshold:.2f} threshold")
    return ssim, psnr


def validate_result(source_format: str, target_format: str, source: Path, output: Path, pair_dir: Path, imagemagick: str, ffprobe: str) -> tuple[list[str], dict[str, float]]:
    failures: list[str] = []
    metrics: dict[str, float] = {}
    if not output.exists() or output.stat().st_size == 0:
        return ["output is missing or empty"], metrics

    document_targets = {"pdf", "docx", "txt", "html", "odt", "xlsx", "ods", "pptx", "odp"}
    if target_format in document_targets:
        text = extract_text(output)
        if MARKER not in text:
            failures.append("expected text marker was not preserved")

    image_targets = {"png", "jpg", "webp", "gif", "tiff", "bmp"}
    if target_format in image_targets:
        try:
            width, height = image_dimensions(output, imagemagick)
            if (width, height) != (320, 180) and source_format != "pdf":
                failures.append(f"unexpected dimensions: {width}x{height}")
            if source_format in {"png", "jpg", "svg", "gif", "tiff"} or source_format == "pdf":
                ssim, psnr = visual_quality(source, output, source_format, target_format, pair_dir, imagemagick)
                metrics["ssim"] = ssim
                metrics["psnr_db"] = psnr
        except Exception as error:
            failures.append(f"image validation failed: {error}")

    media_targets = {"mp4", "webm", "mp3", "wav"}
    if target_format in media_targets:
        try:
            info = media_info(output, ffprobe)
            duration = float(info.get("format", {}).get("duration", 0))
            if abs(duration - EXPECTED_DURATION_SECONDS) > 0.5:
                failures.append(f"duration outside tolerance: {duration:.3f}s")
            if not info.get("streams"):
                failures.append("no media streams detected")
        except Exception as error:
            failures.append(f"media validation failed: {error}")

    if source_format == "pdf" and target_format in {"png", "jpg"}:
        try:
            width, height = image_dimensions(output, imagemagick)
            if width <= 0 or height <= 0:
                failures.append("PDF first-page raster has invalid dimensions")
        except Exception as error:
            failures.append(f"PDF raster validation failed: {error}")
    return failures, metrics


def run_matrix(root: Path) -> int:
    fixtures = generate_fixtures(root)
    results: list[dict[str, object]] = []
    registry = LocalConverterRegistry()
    imagemagick = shutil.which("magick") or binary("convert")
    ffprobe = binary("ffprobe")

    for index, (source_format, target_format) in enumerate(MATRIX, start=1):
        pair_dir = root / "results" / f"{index:02d}-{source_format}-to-{target_format}"
        pair_dir.mkdir(parents=True)
        source = fixtures[source_format]
        result: dict[str, object] = {"source": source_format, "target": target_format, "status": "FAILED"}
        try:
            output = registry.convert(source_format, target_format, source, pair_dir, 900)
            result["output"] = output.name
            errors, metrics = validate_result(source_format, target_format, source, output, pair_dir, imagemagick, ffprobe)
            result["errors"] = errors
            result["metrics"] = metrics
            result["status"] = "PASSED" if not errors else "FAILED"
        except Exception as error:
            result["errors"] = [f"{type(error).__name__}: {error}"]
        results.append(result)
        print(f"[{index:02d}/{len(MATRIX)}] {source_format.upper()} -> {target_format.upper()}: {result['status']}")

    passed = sum(result["status"] == "PASSED" for result in results)
    failed = len(results) - passed
    print(f"\nConversion matrix: {passed}/{len(results)} passed, {failed} failed")
    if failed:
        print("\nQuality report:")
        print(json.dumps([result for result in results if result["status"] == "FAILED"], indent=2, ensure_ascii=False))
    return 0 if failed == 0 else 1


def main() -> int:
    parser = argparse.ArgumentParser(description="Run Morphix's real conversion matrix in a temporary workspace.")
    parser.add_argument("--workspace", type=Path, default=None)
    args = parser.parse_args()

    if args.workspace:
        args.workspace.mkdir(parents=True, exist_ok=True)
        return run_matrix(args.workspace)
    with tempfile.TemporaryDirectory(prefix="morphix-conversions-") as temporary:
        return run_matrix(Path(temporary))


if __name__ == "__main__":
    sys.exit(main())
